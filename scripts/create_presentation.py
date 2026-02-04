"""
Script de generation de la presentation PowerPoint pour L'IA Pero
Selon les criteres du Guide RNCP40875 - Bloc 2
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Couleurs du theme Speakeasy
NOIR = RGBColor(13, 13, 13)
OR = RGBColor(212, 175, 55)
OR_CLAIR = RGBColor(255, 215, 0)
CREME = RGBColor(245, 230, 200)
GRIS_OR = RGBColor(168, 153, 104)

def set_slide_background(slide, color=NOIR):
    """Definir le fond de la slide en noir"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """Slide 1: Page de garde"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Logo martini (texte stylise)
    logo = slide.shapes.add_textbox(Inches(4.5), Inches(0.8), Inches(1), Inches(0.8))
    tf = logo.text_frame
    p = tf.paragraphs[0]
    p.text = "Y"
    p.font.size = Pt(72)
    p.font.color.rgb = OR
    p.alignment = PP_ALIGN.CENTER

    # Titre principal
    title = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "L'IA PERO"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = OR
    p.alignment = PP_ALIGN.CENTER

    # Sous-titre
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.6))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Agent Intelligent Semantique & Generatif"
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = GRIS_OR
    p.alignment = PP_ALIGN.CENTER

    # Encadre RNCP
    box = slide.shapes.add_textbox(Inches(2.5), Inches(4), Inches(5), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "SOUTENANCE PROJET - BLOC 2"
    p.font.size = Pt(18)
    p.font.color.rgb = CREME
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "RNCP40875 - Expert en Ingenierie de Donnees"
    p2.font.size = Pt(16)
    p2.font.color.rgb = CREME
    p2.alignment = PP_ALIGN.CENTER

    # Auteurs
    authors = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    tf = authors.text_frame
    p = tf.paragraphs[0]
    p.text = "Adam BELOUCIF & Amina MEDJDOUB"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Mastere Data Engineering et IA - EFREI Paris"
    p2.font.size = Pt(14)
    p2.font.color.rgb = GRIS_OR
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "Tutrice: Sarah MALAEB | 2025-2026"
    p3.font.size = Pt(12)
    p3.font.color.rgb = GRIS_OR
    p3.alignment = PP_ALIGN.CENTER

def add_context_slide(prs):
    """Slide 2: Problematique et Solution"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "CONTEXTE & PROBLEMATIQUE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = OR

    # Ligne decorative
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()

    # Problematique
    prob_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(0.5))
    tf = prob_title.text_frame
    p = tf.paragraphs[0]
    p.text = "LA PROBLEMATIQUE"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = OR

    prob_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(2))
    tf = prob_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Les moteurs de recherche de recettes classiques (filtres par ingredient) sont:"
    p.font.size = Pt(14)
    p.font.color.rgb = CREME

    points = [
        "- Froids et limites",
        "- Ne capturent pas l'intention utilisateur",
        "- Ignorent le contexte emotionnel",
        "Ex: 'Je veux quelque chose de frais pour l'ete'"
    ]
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(13)
        p.font.color.rgb = CREME
        p.space_before = Pt(6)

    # Solution
    sol_title = slide.shapes.add_textbox(Inches(5), Inches(1.3), Inches(4.5), Inches(0.5))
    tf = sol_title.text_frame
    p = tf.paragraphs[0]
    p.text = "LA SOLUTION: L'IA PERO"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = OR

    sol_content = slide.shapes.add_textbox(Inches(5), Inches(1.8), Inches(4.5), Inches(2.5))
    tf = sol_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Un Barman IA immersif combinant:"
    p.font.size = Pt(14)
    p.font.color.rgb = CREME

    solutions = [
        ("Recherche Semantique (SBERT)", "Comprendre le contexte et l'intention"),
        ("IA Generative (Gemini)", "Creativite et personnalisation"),
        ("Architecture RAG", "Fiabilite et controle des reponses"),
        ("Interface Speakeasy", "Experience utilisateur immersive")
    ]
    for title_text, desc in solutions:
        p = tf.add_paragraph()
        p.text = f"• {title_text}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = OR_CLAIR
        p.space_before = Pt(8)
        p = tf.add_paragraph()
        p.text = f"  {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = CREME

    # Competence RNCP
    comp_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.8))
    tf = comp_box.text_frame
    p = tf.paragraphs[0]
    p.text = "VALIDATION COMPETENCE C4.1"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = OR
    p = tf.add_paragraph()
    p.text = "Definition d'une strategie d'integration IA alignee avec les besoins utilisateurs et l'ecosysteme (Mode hybride Local/Cloud)"
    p.font.size = Pt(11)
    p.font.color.rgb = CREME

def add_user_needs_slide(prs):
    """Slide 3: Analyse du besoin utilisateur"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "ANALYSE DU BESOIN UTILISATEUR"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = OR

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()

    # Persona
    persona_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.5))
    tf = persona_title.text_frame
    p = tf.paragraphs[0]
    p.text = "PERSONA CIBLE"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    persona_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(2))
    tf = persona_content.text_frame
    tf.word_wrap = True
    personas = [
        "• Amateurs de cocktails (25-45 ans)",
        "• Bartenders en quete d'inspiration",
        "• Organisateurs d'evenements",
        "• Curieux de mixologie"
    ]
    for i, persona in enumerate(personas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = persona
        p.font.size = Pt(14)
        p.font.color.rgb = CREME
        p.space_before = Pt(8)

    # Scenarios
    scenario_title = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(0.5))
    tf = scenario_title.text_frame
    p = tf.paragraphs[0]
    p.text = "SCENARIOS D'USAGE"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    scenario_content = slide.shapes.add_textbox(Inches(5), Inches(1.6), Inches(4.5), Inches(2.5))
    tf = scenario_content.text_frame
    tf.word_wrap = True
    scenarios = [
        ("Recherche libre", "'Je veux un cocktail fruite et frais'"),
        ("Filtrage structure", "Budget, Niveau, Temps de preparation"),
        ("Decouverte", "Bouton 'Surprends-moi!'"),
        ("Historique", "Revisiter les creations precedentes")
    ]
    for i, (title_text, desc) in enumerate(scenarios):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {title_text}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = OR_CLAIR
        p.space_before = Pt(8)
        p = tf.add_paragraph()
        p.text = f"  {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = CREME

    # Questionnaire Hybride
    hybrid_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.4))
    tf = hybrid_title.text_frame
    p = tf.paragraphs[0]
    p.text = "QUESTIONNAIRE HYBRIDE (Exigence EF1.1) - IMPLEMENTE"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OR

    hybrid_content = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.2))
    tf = hybrid_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Texte libre + Budget dropdown + 7 Sliders Likert (1-5): Douceur, Acidite, Amertume, Force, Fraicheur, Complexite, Exotisme"
    p.font.size = Pt(13)
    p.font.color.rgb = CREME

def add_methodology_slide(prs):
    """Slide 4: Methodologie et Gestion de projet"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "METHODOLOGIE & ORGANISATION"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = OR

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()

    # Methode Agile
    agile_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.5))
    tf = agile_title.text_frame
    p = tf.paragraphs[0]
    p.text = "APPROCHE AGILE"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    agile_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(2.5))
    tf = agile_content.text_frame
    tf.word_wrap = True
    iterations = [
        "Sprint 1: MVP Backend (Guardrail + Cache)",
        "Sprint 2: Interface Speakeasy (UI/UX)",
        "Sprint 3: Integration GenAI (Gemini)",
        "Sprint 4: Optimisations + Tests E2E"
    ]
    for i, iteration in enumerate(iterations):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {iteration}"
        p.font.size = Pt(13)
        p.font.color.rgb = CREME
        p.space_before = Pt(10)

    # Outils
    tools_title = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(0.5))
    tf = tools_title.text_frame
    p = tf.paragraphs[0]
    p.text = "OUTILS COLLABORATIFS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    tools_content = slide.shapes.add_textbox(Inches(5), Inches(1.6), Inches(4.5), Inches(2.5))
    tf = tools_content.text_frame
    tf.word_wrap = True
    tools = [
        ("GitHub", "Versioning, branches features"),
        ("VS Code", "IDE Python avec extensions"),
        ("Streamlit Cloud", "Deploiement & tests"),
        ("Pytest + Playwright", "Tests E2E automatises")
    ]
    for i, (tool, desc) in enumerate(tools):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {tool}: {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = CREME
        p.space_before = Pt(10)

    # Structure projet
    struct_title = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.4))
    tf = struct_title.text_frame
    p = tf.paragraphs[0]
    p.text = "STRUCTURE DU PROJET (~3,300 lignes Python)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OR

    struct_content = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9), Inches(1))
    tf = struct_content.text_frame
    p = tf.paragraphs[0]
    p.text = "ia-pero/ -> src/ (app.py, backend.py, scoring.py) | data/ (CSV, JSON) | tests/ (19 tests Pytest) | scripts/"
    p.font.size = Pt(12)
    p.font.color.rgb = CREME

def add_data_slide(prs):
    """Slide 5: Referentiel de donnees"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "REFERENTIEL DE DONNEES"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = OR

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()

    # Dataset
    dataset_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.5))
    tf = dataset_title.text_frame
    p = tf.paragraphs[0]
    p.text = "DATASET PROPRIETAIRE"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    dataset_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(2))
    tf = dataset_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "600 cocktails generes (src/generate_data.py)"
    p.font.size = Pt(14)
    p.font.color.rgb = CREME
    p = tf.add_paragraph()
    p.text = "+ 1000 cocktails Kaggle enrichis"
    p.font.size = Pt(14)
    p.font.color.rgb = CREME
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = "= 1600+ cocktails disponibles"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = OR_CLAIR
    p.space_before = Pt(12)

    # Structure
    struct_title = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(0.5))
    tf = struct_title.text_frame
    p = tf.paragraphs[0]
    p.text = "STRUCTURE DU DATASET"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    struct_content = slide.shapes.add_textbox(Inches(5), Inches(1.6), Inches(4.5), Inches(2.5))
    tf = struct_content.text_frame
    tf.word_wrap = True
    fields = [
        ("name", "Nom du cocktail"),
        ("description_semantique", "Texte riche pour embeddings SBERT"),
        ("ingredients", "Liste des ingredients"),
        ("taste_profile", "JSON 7 dimensions (radar chart)"),
        ("source", "'generated' ou 'kaggle'")
    ]
    for i, (field, desc) in enumerate(fields):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {field}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = OR_CLAIR
        p.space_before = Pt(6)
        p = tf.add_paragraph()
        p.text = f"  {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = CREME

    # Competence
    comp_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
    tf = comp_box.text_frame
    p = tf.paragraphs[0]
    p.text = "VALIDATION COMPETENCE C3.1"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = OR
    p = tf.add_paragraph()
    p.text = "Preparer, nettoyer et structurer les donnees pour assurer une qualite optimale pour l'IA"
    p.font.size = Pt(11)
    p.font.color.rgb = CREME

def add_pipeline_slide(prs):
    """Slide 6: Pipeline IA et Architecture"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "PIPELINE IA & ARCHITECTURE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = OR

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()

    # Pipeline steps
    steps = [
        ("1. INPUT", "Questionnaire hybride\n(texte + budget)"),
        ("2. GUARDRAIL", "Check SBERT\nsimilarite > 0.30"),
        ("3. CACHE", "MD5 Hash\nLookup JSON"),
        ("4. EMBEDDINGS", "SBERT encode\n384 dimensions"),
        ("5. SCORING", "Similarite Cosinus\nTop-N ranking"),
        ("6. RAG/GENAI", "Gemini API\n(si pas en cache)"),
        ("7. OUTPUT", "Recipe + Radar\n+ Export PDF")
    ]

    x_start = 0.3
    box_width = 1.25
    box_height = 1.1
    y_pos = 1.4

    for i, (step_title, step_desc) in enumerate(steps):
        x_pos = x_start + i * (box_width + 0.08)

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(y_pos), Inches(box_width), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(26, 26, 26)
        box.line.color.rgb = OR
        box.line.width = Pt(1)

        # Title
        title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.05), Inches(box_width), Inches(0.3))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = OR
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.35), Inches(box_width), Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(8)
        p.font.color.rgb = CREME
        p.alignment = PP_ALIGN.CENTER

        # Arrow (except last)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_textbox(Inches(x_pos + box_width), Inches(y_pos + 0.4), Inches(0.15), Inches(0.3))
            tf = arrow.text_frame
            p = tf.paragraphs[0]
            p.text = ">"
            p.font.size = Pt(14)
            p.font.color.rgb = OR

    # Details techniques
    tech_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(4.2), Inches(0.4))
    tf = tech_title.text_frame
    p = tf.paragraphs[0]
    p.text = "SBERT LOCAL (all-MiniLM-L6-v2)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OR

    tech_content = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(4.2), Inches(1.5))
    tf = tech_content.text_frame
    tf.word_wrap = True
    points = [
        "• Cout: 0 EUR (modele local)",
        "• Latence: ~50ms vs 500ms API Cloud",
        "• Confidentialite: Donnees on-premise",
        "• Embeddings: 384 dimensions"
    ]
    for i, point in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = point
        p.font.size = Pt(12)
        p.font.color.rgb = CREME
        p.space_before = Pt(6)

    # GenAI
    genai_title = slide.shapes.add_textbox(Inches(5), Inches(2.8), Inches(4.5), Inches(0.4))
    tf = genai_title.text_frame
    p = tf.paragraphs[0]
    p.text = "GOOGLE GEMINI (RAG)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OR

    genai_content = slide.shapes.add_textbox(Inches(5), Inches(3.2), Inches(4.5), Inches(1.5))
    tf = genai_content.text_frame
    tf.word_wrap = True
    points = [
        "• Modele: gemini-2.5-flash-lite",
        "• Fallback: 5 modeles en cascade",
        "• Persona: 'Barman annees 1920'",
        "• Output: JSON structure valide"
    ]
    for i, point in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = point
        p.font.size = Pt(12)
        p.font.color.rgb = CREME
        p.space_before = Pt(6)

def add_implementation_slide(prs):
    """Slide 7: Implementation Technique"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "IMPLEMENTATION TECHNIQUE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = OR

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()

    # Technologies
    tech_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.4))
    tf = tech_title.text_frame
    p = tf.paragraphs[0]
    p.text = "STACK TECHNIQUE"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    tech_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(2.2))
    tf = tech_content.text_frame
    tf.word_wrap = True
    techs = [
        ("Python 3.11", "Langage principal"),
        ("Streamlit", "Interface web reactive"),
        ("Sentence-Transformers", "Embeddings SBERT"),
        ("Google GenAI", "API Gemini"),
        ("Plotly", "Graphiques interactifs"),
        ("Pandas/NumPy", "Manipulation donnees")
    ]
    for i, (tech, desc) in enumerate(techs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {tech}: {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = CREME
        p.space_before = Pt(6)

    # Scoring
    score_title = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = score_title.text_frame
    p = tf.paragraphs[0]
    p.text = "SCORING RNCP (EF3.1)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    score_content = slide.shapes.add_textbox(Inches(5), Inches(1.6), Inches(4.5), Inches(2.2))
    tf = score_content.text_frame
    tf.word_wrap = True
    scores = [
        "• Coverage Score = SWi*Si / SWi",
        "• Similarite Cosinus par bloc",
        "• 7 dimensions gustatives ponderees",
        "• Preferences Likert (1-5) integrees",
        "• Affichage score % + barres progression"
    ]
    for i, score in enumerate(scores):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = score
        p.font.size = Pt(12)
        p.font.color.rgb = CREME
        p.space_before = Pt(6)

    # Gouvernance
    gov_title = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.4))
    tf = gov_title.text_frame
    p = tf.paragraphs[0]
    p.text = "GOUVERNANCE & RESPONSABILISATION"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OR

    gov_content = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9), Inches(1))
    tf = gov_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• Guardrail semantique anti-injection | Cache JSON pour audit | Fallback local si API down | Analytics loguees | Mention legale alcool"
    p.font.size = Pt(11)
    p.font.color.rgb = CREME

def add_dataviz_slide(prs):
    """Slide 8: DataViz et UX"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "DATAVIZ & EXPERIENCE UTILISATEUR"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = OR

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()

    # Radar Chart
    radar_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.4))
    tf = radar_title.text_frame
    p = tf.paragraphs[0]
    p.text = "GRAPHIQUE RADAR (7 AXES)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    radar_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(2.5))
    tf = radar_content.text_frame
    tf.word_wrap = True
    axes = [
        "• Douceur (1.5 - 5.0)",
        "• Acidite (1.5 - 5.0)",
        "• Amertume (1.5 - 5.0)",
        "• Force alcool (1.5 - 5.0)",
        "• Fraicheur (1.5 - 5.0)",
        "• Prix (1.5 - 5.0)",
        "• Qualite (1.5 - 5.0)"
    ]
    for i, axe in enumerate(axes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = axe
        p.font.size = Pt(12)
        p.font.color.rgb = CREME
        p.space_before = Pt(4)

    # UI Features
    ui_title = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = ui_title.text_frame
    p = tf.paragraphs[0]
    p.text = "INTERFACE STREAMLIT"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    ui_content = slide.shapes.add_textbox(Inches(5), Inches(1.6), Inches(4.5), Inches(2.5))
    tf = ui_content.text_frame
    tf.word_wrap = True
    features = [
        "• Theme Speakeasy annees 1920",
        "• 7 Sliders Likert pour profil gustatif",
        "• Coverage Score RNCP avec barres",
        "• Plan de Decouverte personnalise",
        "• Bio du profil gustatif",
        "• Enrichissement auto requetes courtes",
        "• Export recette + 'Surprends-moi!'"
    ]
    for i, feat in enumerate(features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = feat
        p.font.size = Pt(12)
        p.font.color.rgb = CREME
        p.space_before = Pt(4)

    # Competence
    comp_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
    tf = comp_box.text_frame
    p = tf.paragraphs[0]
    p.text = "VALIDATION COMPETENCE C3.2"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = OR
    p = tf.add_paragraph()
    p.text = "Elaborer une communication infographique visuelle via des tableaux de bord interactifs"
    p.font.size = Pt(11)
    p.font.color.rgb = CREME

def add_results_slide(prs):
    """Slide 9: Resultats et Performance"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "RESULTATS & PERFORMANCE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = OR

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()

    # Metriques
    metrics = [
        ("95%", "Taux de rejet correct\n(Guardrail SBERT)"),
        ("~50ms", "Temps reponse\n(avec cache)"),
        ("40-60x", "Speedup recherche\n(embeddings caches)"),
        ("100%", "Validite JSON\n(prompt strict)")
    ]

    x_start = 0.5
    box_width = 2.1

    for i, (value, label) in enumerate(metrics):
        x_pos = x_start + i * (box_width + 0.2)

        # Value
        val_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.3), Inches(box_width), Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = OR_CLAIR
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(x_pos), Inches(2.1), Inches(box_width), Inches(0.8))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = CREME
        p.alignment = PP_ALIGN.CENTER

    # Tests
    test_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(4.2), Inches(0.4))
    tf = test_title.text_frame
    p = tf.paragraphs[0]
    p.text = "TESTS AUTOMATISES"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OR

    test_content = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(4.2), Inches(1.5))
    tf = test_content.text_frame
    tf.word_wrap = True
    tests = [
        "• 19 tests automatises (Pytest)",
        "• Playwright E2E: Flow complet (3)",
        "• Scoring module: 16 tests unitaires",
        "• Integration Kaggle: 5/5 passes"
    ]
    for i, test in enumerate(tests):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = test
        p.font.size = Pt(12)
        p.font.color.rgb = CREME
        p.space_before = Pt(6)

    # Green AI
    green_title = slide.shapes.add_textbox(Inches(5), Inches(3.2), Inches(4.5), Inches(0.4))
    tf = green_title.text_frame
    p = tf.paragraphs[0]
    p.text = "GREEN AI & OPTIMISATION"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OR

    green_content = slide.shapes.add_textbox(Inches(5), Inches(3.6), Inches(4.5), Inches(1.5))
    tf = green_content.text_frame
    tf.word_wrap = True
    greens = [
        "• MiniLM local (80 Mo) vs LLM geant",
        "• Cache JSON: 0 appel API redondant",
        "• Empreinte memoire: ~180 Mo",
        "• CPU standard (pas de GPU requis)"
    ]
    for i, green in enumerate(greens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = green
        p.font.size = Pt(12)
        p.font.color.rgb = CREME
        p.space_before = Pt(6)

def add_limits_slide(prs):
    """Slide 10: Limites et Ameliorations"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "LIMITES & PISTES D'AMELIORATION"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = OR

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()

    # Limites
    limits_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.4))
    tf = limits_title.text_frame
    p = tf.paragraphs[0]
    p.text = "LIMITES IDENTIFIEES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    limits_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(2.5))
    tf = limits_content.text_frame
    tf.word_wrap = True
    limits = [
        "• SBERT anglais: Performances reduites en francais",
        "• Dataset limite: 1600 cocktails",
        "• Dependance API Gemini: Quota 15 req/min",
        "• Variabilite des similarites selon formulation",
        "• Pas de persistance profil (session only)"
    ]
    for i, limit in enumerate(limits):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = limit
        p.font.size = Pt(12)
        p.font.color.rgb = CREME
        p.space_before = Pt(8)

    # Ameliorations
    improve_title = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(0.4))
    tf = improve_title.text_frame
    p = tf.paragraphs[0]
    p.text = "AMELIORATIONS FUTURES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OR

    improve_content = slide.shapes.add_textbox(Inches(5), Inches(1.6), Inches(4.5), Inches(2.5))
    tf = improve_content.text_frame
    tf.word_wrap = True
    improvements = [
        "• Fine-tuning SBERT sur corpus cocktails",
        "• Base de donnees vectorielle (Pinecone/Chroma)",
        "• Chatbot conversationnel guide",
        "• Profils utilisateurs persistants",
        "• API REST pour integration externe"
    ]
    for i, improve in enumerate(improvements):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = improve
        p.font.size = Pt(12)
        p.font.color.rgb = CREME
        p.space_before = Pt(8)

def add_conclusion_slide(prs):
    """Slide 11: Conclusion"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Titre
    title = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(6), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "CONCLUSION"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = OR
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(8), Inches(0.6))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Projet validant l'ensemble du Bloc 2 RNCP40875"
    p.font.size = Pt(20)
    p.font.color.rgb = GRIS_OR
    p.alignment = PP_ALIGN.CENTER

    # Competences boxes
    competences = [
        ("DATA ENG.", "Pipeline de donnees\nCSV, Nettoyage\n600+ cocktails"),
        ("MACHINE LEARNING", "NLP Semantique\nEmbeddings SBERT\nSimilarite cosinus"),
        ("GENAI", "Integration LLM\nArchitecture RAG\nPrompt Engineering"),
        ("INDUSTRIALISATION", "UI Streamlit\nTests E2E\nCache & Logging")
    ]

    x_start = 0.4
    box_width = 2.2
    y_pos = 2.2

    for i, (comp_title, comp_desc) in enumerate(competences):
        x_pos = x_start + i * (box_width + 0.15)

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(y_pos), Inches(box_width), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(26, 26, 26)
        box.line.color.rgb = OR
        box.line.width = Pt(1)

        # Title
        title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.1), Inches(box_width), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = comp_title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = OR
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(y_pos + 0.5), Inches(box_width - 0.2), Inches(1))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = comp_desc
        p.font.size = Pt(10)
        p.font.color.rgb = CREME
        p.alignment = PP_ALIGN.CENTER

    # Demo call
    demo = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(0.6))
    tf = demo.text_frame
    p = tf.paragraphs[0]
    p.text = "Place a la demonstration!"
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = OR_CLAIR
    p.alignment = PP_ALIGN.CENTER

    # Martini icon
    icon = slide.shapes.add_textbox(Inches(4.5), Inches(4.8), Inches(1), Inches(0.8))
    tf = icon.text_frame
    p = tf.paragraphs[0]
    p.text = "Y"
    p.font.size = Pt(48)
    p.font.color.rgb = OR
    p.alignment = PP_ALIGN.CENTER

def add_demo_slide(prs):
    """Slide 12: Demo Live"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "DEMONSTRATION LIVE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = OR

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()

    # Scenarios
    scenarios = [
        ("1. PROFIL LIKERT", "Ajuster les 7 sliders de preferences\n-> Enrichissement auto des requetes courtes (EF4.1)"),
        ("2. COVERAGE SCORE", "Requete: 'mojito' -> Score RNCP %\n-> Formule SWi*Si/SWi + barres progression (EF3.1)"),
        ("3. PLAN DECOUVERTE", "Generation cocktail\n-> Plan de progression + Bio profil (EF4.2/EF4.3)"),
        ("4. GENERATION GENAI", "Requete: 'cocktail mysterieux'\n-> Recipe + Radar + Scoring + Export")
    ]

    y_start = 1.3

    for i, (scenario_title, scenario_desc) in enumerate(scenarios):
        y_pos = y_start + i * 1.1

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(3), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = scenario_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = OR

        # Description
        desc_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_pos), Inches(6), Inches(0.9))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = scenario_desc
        p.font.size = Pt(13)
        p.font.color.rgb = CREME

    # Note
    note = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.text = "Commande: streamlit run src/app.py | URL: http://localhost:8501"
    p.font.size = Pt(14)
    p.font.color.rgb = GRIS_OR
    p.alignment = PP_ALIGN.CENTER

def main():
    """Generer la presentation complete"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # Ajouter toutes les slides
    add_title_slide(prs)           # 1. Page de garde
    add_context_slide(prs)         # 2. Contexte & Problematique
    add_user_needs_slide(prs)      # 3. Analyse besoin utilisateur
    add_methodology_slide(prs)     # 4. Methodologie
    add_data_slide(prs)            # 5. Referentiel de donnees
    add_pipeline_slide(prs)        # 6. Pipeline IA
    add_implementation_slide(prs)  # 7. Implementation technique
    add_dataviz_slide(prs)         # 8. DataViz & UX
    add_results_slide(prs)         # 9. Resultats & Performance
    add_limits_slide(prs)          # 10. Limites & Ameliorations
    add_conclusion_slide(prs)      # 11. Conclusion
    add_demo_slide(prs)            # 12. Demo Live

    # Sauvegarder
    output_path = "L'IA Pero - Soutenance RNCP Bloc 2 v2.pptx"
    prs.save(output_path)
    print(f"Presentation generee: {output_path}")
    print(f"Nombre de slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    main()
